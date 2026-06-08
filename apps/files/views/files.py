# -*- coding: utf-8 -*-
"""
apps.files.views.files - 文件管理视图

提供文件上传、预览、视频列表与详情等接口。
"""

import logging
import mimetypes
import os
import uuid

from pathlib import Path

from django.conf import settings
from django.http import FileResponse, Http404
from django.shortcuts import get_object_or_404
from django.utils import timezone

from rest_framework import generics, status
from rest_framework.decorators import api_view, permission_classes
from rest_framework.parsers import FormParser, MultiPartParser
from rest_framework.permissions import AllowAny, IsAuthenticated
from rest_framework.response import Response
from rest_framework.views import APIView

from apps.files.models import Video
from apps.files.serializers import FileUploadSerializer, VideoSerializer

logger = logging.getLogger(__name__)

# ============================================================================
# 常量
# ============================================================================

# 上传文件限制：单个文件最大 200 MB
MAX_UPLOAD_SIZE = 200 * 1024 * 1024

# 允许的文件扩展名映射
ALLOWED_EXTENSIONS = {
    'image': {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'},
    'pdf': {'.pdf'},
    'video': {'.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv', '.webm'},
    'ppt': {'.pptx', '.ppt'},
}

# 文件大小展示单位
SIZE_UNITS = ['B', 'KB', 'MB', 'GB']


# ============================================================================
# 辅助函数
# ============================================================================


def _format_file_size(size_bytes: int) -> str:
    """将字节转换为人类可读的文件大小"""
    if size_bytes <= 0:
        return '0 B'
    unit_index = 0
    size = float(size_bytes)
    while size >= 1024 and unit_index < len(SIZE_UNITS) - 1:
        size /= 1024
        unit_index += 1
    return f'{size:.2f} {SIZE_UNITS[unit_index]}'


def _detect_file_type(filename: str) -> str:
    """根据文件扩展名推断文件类型"""
    ext = Path(filename).suffix.lower()
    for ftype, extensions in ALLOWED_EXTENSIONS.items():
        if ext in extensions:
            return ftype
    return 'other'


def _validate_file_extension(filename: str, expected_type: str) -> bool:
    """校验文件扩展名是否与声明的类型匹配"""
    ext = Path(filename).suffix.lower()
    allowed = ALLOWED_EXTENSIONS.get(expected_type, set())
    return ext in allowed


# ============================================================================
# 视图
# ============================================================================


class FileUploadView(APIView):
    """
    POST /api/v1/files/upload/

    文件上传接口（需登录）

    支持的文件类型：
      - image: jpg, jpeg, png, gif, bmp, webp
      - pdf: pdf
      - video: mp4, avi, mov, wmv, flv, mkv, webm
      - ppt: pptx, ppt

    请求体（multipart/form-data）：
      - file: 文件二进制
      - type: 文件类型（image/pdf/video/ppt/other）

    响应：
      {
        "code": 0,
        "message": "上传成功",
        "data": {
          "url": "/media/uploads/xxx.ext",
          "fileName": "xxx.ext",
          "fileSize": "1.23 MB",
          "fileType": "image",
          "extension": ".jpg"
        }
      }
    """

    parser_classes = [MultiPartParser, FormParser]
    permission_classes = [IsAuthenticated]
    serializer_class = FileUploadSerializer

    def post(self, request, *args, **kwargs):
        serializer = self.serializer_class(data=request.data)
        serializer.is_valid(raise_exception=True)

        uploaded_file = serializer.validated_data['file']
        file_type = serializer.validated_data.get('type', 'other')
        experiment_id = serializer.validated_data.get('experimentId', '')

        # --- 文件大小校验 ---
        if uploaded_file.size > MAX_UPLOAD_SIZE:
            return Response(
                {
                    'code': 'FILE_TOO_LARGE',
                    'message': f'文件大小不能超过 200MB',
                    'details': {},
                },
                status=status.HTTP_400_BAD_REQUEST,
            )

        # --- 原始文件名与扩展名 ---
        original_name = uploaded_file.name
        ext = Path(original_name).suffix.lower()

        # --- 类型匹配校验 ---
        if file_type != 'other' and not _validate_file_extension(original_name, file_type):
            allowed_str = ', '.join(sorted(ALLOWED_EXTENSIONS.get(file_type, set())))
            return Response(
                {
                    'code': 'INVALID_FILE_TYPE',
                    'message': f'文件类型 "{file_type}" 不支持扩展名 "{ext}"，'
                               f'仅支持: {allowed_str}',
                    'details': {},
                },
                status=status.HTTP_400_BAD_REQUEST,
            )

        # --- 生成唯一文件名，避免重名覆盖 ---
        unique_name = f'{uuid.uuid4().hex}{ext}'

        # 按类型划分子目录：media/uploads/{type}/
        upload_subdir = Path('uploads') / file_type
        upload_dir = settings.MEDIA_ROOT / upload_subdir
        upload_dir.mkdir(parents=True, exist_ok=True)

        # 保存文件
        dest_path = upload_dir / unique_name
        with open(dest_path, 'wb') as f:
            for chunk in uploaded_file.chunks():
                f.write(chunk)

        # --- 构建返回 URL ---
        file_url = f'{settings.MEDIA_URL}{upload_subdir.as_posix()}/{unique_name}'

        # --- 如果是实验报告（传入了 experimentId），关联到 tb_experiment_score ---
        if experiment_id:
            from django.utils.timezone import now
            from django.db import connection

            user_id = request.user.id
            with connection.cursor() as cur:
                # 查找是否已有该实验的评分记录
                cur.execute(
                    "SELECT id FROM tb_experiment_score WHERE user_id=%s AND experiment_id=%s",
                    [user_id, experiment_id]
                )
                exist_row = cur.fetchone()
                if exist_row:
                    # 更新 pdf_path
                    cur.execute(
                        "UPDATE tb_experiment_score SET pdf_path=%s, update_time=%s WHERE id=%s",
                        [str(dest_path), now(), exist_row[0]]
                    )
                else:
                    # 创建新评分记录
                    new_id = uuid.uuid4().hex[:32]
                    cur.execute(
                        "INSERT INTO tb_experiment_score (id, user_id, experiment_id, pdf_path, create_time, update_time) VALUES (%s,%s,%s,%s,%s,%s)",
                        [new_id, user_id, experiment_id, str(dest_path), now(), now()]
                    )

        logger.info(
            'File uploaded: user=%s, type=%s, name=%s, size=%d, path=%s',
            request.user.id,
            file_type,
            original_name,
            uploaded_file.size,
            dest_path,
        )

        return Response(
            {
                'code': 0,
                'message': '上传成功',
                'data': {
                    'url': file_url,
                    'fileName': original_name,
                    'fileSize': _format_file_size(uploaded_file.size),
                    'fileType': file_type,
                    'extension': ext,
                },
            },
            status=status.HTTP_201_CREATED,
        )


class FilePreviewView(APIView):
    """
    GET /api/v1/files/preview/<path:file_path>/

    文件预览接口（匿名可访问）

    支持的处理方式：
      - 图片：直接返回文件流
      - PDF：返回第一页图片预览
      - PPT：将第一页转换为图片后返回
      - 视频：返回文件流（浏览器直接播放）
      - 其他：返回文件流

    注意：file_path 是相对于 media/uploads/ 的路径，例如：
      - image/abc123.jpg
      - pdf/def456.pdf
      - ppt/ghi789.pptx
    """

    permission_classes = [AllowAny]

    def get(self, request, file_path, *args, **kwargs):
        # 确保路径安全，防止目录穿越
        safe_path = Path(file_path).as_posix()
        if '..' in safe_path:
            raise Http404('Invalid file path')

        full_path = settings.MEDIA_ROOT / 'uploads' / safe_path
        if not full_path.exists() or not full_path.is_file():
            raise Http404('文件不存在')

        # 推断文件类型
        ext = full_path.suffix.lower()
        file_type = _detect_file_type(full_path.name)
        mime_type, _ = mimetypes.guess_type(str(full_path))

        # --- 图片：直接返回 ---
        if file_type == 'image':
            return FileResponse(
                open(str(full_path), 'rb'),
                content_type=mime_type or 'image/jpeg',
            )

        # --- PDF：转图片预览 ---
        if file_type == 'pdf':
            return self._pdf_preview(full_path)

        # --- PPT：第一页转图片预览 ---
        if file_type == 'ppt':
            return self._ppt_preview(full_path)

        # --- 其他（视频等）：直接返回文件流 ---
        return FileResponse(
            open(str(full_path), 'rb'),
            content_type=mime_type or 'application/octet-stream',
        )

    def _pdf_preview(self, pdf_path: Path):
        """
        PDF 转图片预览

        使用 pdf2image + poppler 将 PDF 第一页转为 PNG 图片返回。
        如果 pdf2image 不可用，回退为直接返回 PDF 文件。
        """
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(
                str(pdf_path),
                first_page=1,
                last_page=1,
                fmt='jpeg',
                dpi=150,
            )
            if not images:
                raise Http404('PDF 解析失败')
            from io import BytesIO
            buf = BytesIO()
            images[0].save(buf, format='JPEG', quality=85)
            buf.seek(0)
            return FileResponse(buf, content_type='image/jpeg')
        except ImportError:
            logger.warning('pdf2image not installed, falling back to raw PDF download')
            return FileResponse(
                open(str(pdf_path), 'rb'),
                content_type='application/pdf',
            )
        except Exception as exc:
            logger.error('PDF preview failed: %s', exc, exc_info=True)
            return FileResponse(
                open(str(pdf_path), 'rb'),
                content_type='application/pdf',
            )

    def _ppt_preview(self, ppt_path: Path):
        """
        PPT 第一页预览

        使用 python-pptx 解析并生成第一页缩略图。
        如果不可用，回退为直接返回文件流。
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from io import BytesIO

            prs = Presentation(str(ppt_path))

            # 创建一个空白预览图（宽高比 16:9），模拟 PPT 第一页
            img_width = 1280
            img_height = 720
            img = Image.new('RGB', (img_width, img_height), (255, 255, 255))
            draw = ImageDraw.Draw(img)

            # 尝试从第一张幻灯片获取文本内容
            if prs.slides:
                slide = prs.slides[0]
                text_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_lines.append(shape.text.strip()[:200])

                # 在图片上绘制文本
                y = 50
                for line in text_lines[:20]:
                    draw.text((50, y), line, fill=(50, 50, 50))
                    y += 30
            else:
                draw.text((50, 50), '(空白演示文稿)', fill=(180, 180, 180))

            # 绘制边框
            draw.rectangle([0, 0, img_width - 1, img_height - 1],
                           outline=(200, 200, 200), width=2)

            buf = BytesIO()
            img.save(buf, format='JPEG', quality=85)
            buf.seek(0)
            return FileResponse(buf, content_type='image/jpeg')

        except ImportError:
            logger.warning('python-pptx not installed, falling back to raw PPT download')
            return FileResponse(
                open(str(ppt_path), 'rb'),
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            )
        except Exception as exc:
            logger.error('PPT preview failed: %s', exc, exc_info=True)
            return FileResponse(
                open(str(ppt_path), 'rb'),
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            )


class VideoListView(generics.ListAPIView):
    """
    GET /api/v1/files/videos/

    视频列表（分页，匿名可访问）
    """

    queryset = Video.objects.all().order_by('-createTime')
    serializer_class = VideoSerializer
    permission_classes = [AllowAny]


class VideoDetailView(generics.RetrieveAPIView):
    """
    GET /api/v1/files/videos/<id>/

    视频详情（匿名可访问）

    返回包含播放 URL 的视频信息。
    """

    queryset = Video.objects.all()
    serializer_class = VideoSerializer
    permission_classes = [AllowAny]


# ============================================================================
# 实验报告
# ============================================================================


@api_view(['GET'])
@permission_classes([AllowAny])
def experiment_report(request, experiment_id):
    """
    GET /api/v1/files/report/<experiment_id>/

    获取当前用户指定实验的实验报告信息。
    如果未登录或没有报告，返回空数据。
    """
    user_id = None
    if hasattr(request, 'user') and request.user and hasattr(request.user, 'id'):
        user_id = request.user.id

    if not user_id:
        return Response({'hasReport': False, 'report': None})

    from django.db import connection
    with connection.cursor() as cur:
        cur.execute(
            "SELECT id, pdf_path, score_sum, create_time, update_time FROM tb_experiment_score WHERE user_id=%s AND experiment_id=%s ORDER BY create_time DESC LIMIT 1",
            [user_id, experiment_id]
        )
        row = cur.fetchone()

    if not row or not row[1]:
        return Response({'hasReport': False, 'report': None})

    pdf_path = row[1]
    # 将文件路径转为 URL
    # 路径类似 media/uploads/pdf/xxx.pdf，转为 /media/uploads/pdf/xxx.pdf
    pdf_url = None
    if pdf_path:
        rel_path = str(pdf_path)
        # 如果含 media/ 前缀，保留相对路径
        if 'media/' in rel_path:
            pdf_url = '/' + rel_path[rel_path.index('media/'):]
        else:
            pdf_url = settings.MEDIA_URL + 'uploads/' + os.path.basename(os.path.dirname(pdf_path)) + '/' + os.path.basename(pdf_path) if os.path.basename(pdf_path) else None
            pdf_url = '/' + pdf_url.lstrip('/') if pdf_url else None

    return Response({
        'hasReport': True,
        'report': {
            'id': row[0],
            'fileUrl': pdf_url,
            'score': float(row[2]) if row[2] else None,
            'createTime': str(row[3])[:19] if row[3] else '',
            'updateTime': str(row[4])[:19] if row[4] else '',
        }
    })
